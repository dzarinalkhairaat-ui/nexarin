from fastapi import APIRouter, File, UploadFile, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse
from pptx import Presentation  # type: ignore
from pptx.util import Inches, Pt  # type: ignore
from pptx.enum.text import MSO_AUTO_SIZE # type: ignore
import fitz  # PyMuPDF
import tempfile
import os
import time

router = APIRouter()

def cleanup_file(filepath: str):
    if filepath and os.path.exists(filepath):
        try:
            os.remove(filepath)
        except:
            pass

@router.post("/convert/pdf-to-pptx")
async def pdf_to_pptx(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Hanya menerima file PDF.")
        
    try:
        doc_bytes = await file.read()
        doc = fitz.open(stream=doc_bytes, filetype="pdf")
        
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        if len(doc) > 0:
            first_page = doc[0]
            width = first_page.rect.width
            height = first_page.rect.height
            prs.slide_width = Inches(width / 72)
            prs.slide_height = Inches(height / 72)
        
        with tempfile.TemporaryDirectory() as temp_dir:
            for i in range(len(doc)):
                page = doc[i]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                try:
                    # Advanced extraction mode: pisahkan teks dan gambar
                    blocks = page.get_text("dict")["blocks"]
                    
                    if not blocks:
                        # Fallback jika tidak ada block sama sekali (misal: scanned PDF)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        img_path = os.path.join(temp_dir, f"page_fb_{i}.png")
                        pix.save(img_path)
                        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                        continue

                    # 1. Ekstrak Gambar terlebih dahulu (agar teks berada di atas gambar)
                    for b_idx, b in enumerate(blocks):
                        if b["type"] == 1:  # image block
                            x0, y0, x1, y1 = b["bbox"]
                            left, top = Inches(x0 / 72), Inches(y0 / 72)
                            w, h = Inches((x1 - x0) / 72), Inches((y1 - y0) / 72)
                            
                            img_bytes = b.get("image")
                            if img_bytes:
                                img_path = os.path.join(temp_dir, f"img_{i}_{b_idx}.png")
                                with open(img_path, "wb") as f:
                                    f.write(img_bytes)
                                try:
                                    slide.shapes.add_picture(img_path, left, top, width=w, height=h)
                                except Exception as img_ex:
                                    print(f"Failed to add picture {b_idx}: {img_ex}")

                    # 2. Ekstrak Teks menjadi TextBox yang bisa diedit
                    for b_idx, b in enumerate(blocks):
                        if b["type"] == 0:  # text block
                            x0, y0, x1, y1 = b["bbox"]
                            left, top = Inches(x0 / 72), Inches(y0 / 72)
                            w, h = Inches((x1 - x0) / 72), Inches((y1 - y0) / 72)
                            
                            txBox = slide.shapes.add_textbox(left, top, w, h)
                            tf = txBox.text_frame
                            tf.word_wrap = True
                            
                            p = tf.paragraphs[0]
                            p.text = ""
                            
                            for line_idx, line in enumerate(b["lines"]):
                                for span in line["spans"]:
                                    run = p.add_run()
                                    run.text = span["text"]
                                    run.font.size = Pt(span["size"])
                                    
                                if line_idx < len(b["lines"]) - 1:
                                    p.add_run().text = "\n"

                except Exception as ex:
                    # Fallback ke rendering gambar penuh jika terjadi error kompleks
                    print(f"Advanced extraction failed for page {i}: {ex}, falling back to image")
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img_path = os.path.join(temp_dir, f"page_fb_err_{i}.png")
                    pix.save(img_path)
                    slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

            # Simpan presentasi ke file sementara
            unique_id = int(time.time())
            output_path = os.path.join(tempfile.gettempdir(), f"output_{unique_id}.pptx")
            prs.save(output_path)
            
        doc.close()
        
        base_name = os.path.splitext(file.filename)[0]
        
        # Bersihkan file PPTX setelah dikirim ke pengguna menggunakan BackgroundTasks
        background_tasks.add_task(cleanup_file, output_path)
        
        return FileResponse(
            output_path, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            filename=f"{base_name}_nexarin.pptx"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Terjadi kesalahan saat mengonversi PDF ke PPTX: {str(e)}")
