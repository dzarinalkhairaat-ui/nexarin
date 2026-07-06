from fastapi import FastAPI, File, UploadFile, HTTPException
from typing import List
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from pdf2docx import Converter
from PIL import Image
import fitz  # PyMuPDF
import io
import os
import tempfile
import time
import random
import subprocess
import zipfile
import pdfplumber  # type: ignore
import pandas as pd  # type: ignore
from pptx import Presentation  # type: ignore
from pptx.util import Inches  # type: ignore

app = FastAPI(title="Nexarin PDF API", description="Microservice for PDF to Word Conversion")

# Allow CORS so Next.js frontend can call this API
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # Change this to your frontend domain in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def read_root():
    return {"message": "Nexarin PDF API is running!"}

@app.post("/convert/pdf-to-word")
async def convert_pdf_to_word(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File must be a PDF")
    
    # Create temporary files
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    input_pdf = os.path.join(temp_dir, f"input_{unique_id}.pdf")
    output_docx = os.path.join(temp_dir, f"output_{unique_id}.docx")
    
    try:
        # Save uploaded PDF to temp file
        with open(input_pdf, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        # Convert using pdf2docx with high accuracy parameters
        cv = Converter(input_pdf)
        kwargs = {
            "shape_min_dimension": 0.5,
            "float_image_ignorable_gap": 0.0,
            "connected_border_tolerance": 1.5,
            "lines_left_aligned_threshold": 0.5,
            "lines_center_aligned_threshold": 1.0,
            "lines_right_aligned_threshold": 0.5,
            "clip_image_res_ratio": 3.0
        }
        cv.convert(output_docx, kwargs=kwargs)
        cv.close()
        
        if not os.path.exists(output_docx):
            raise Exception("DOCX file was not generated")
            
        # Return the file response
        # Note: In a real serverless env, background tasks or similar might be needed to delete the file AFTER sending.
        # FastAPI's FileResponse doesn't delete it automatically, but tempfile directory usually gets cleaned up by OS.
        return FileResponse(
            output_docx, 
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=f"{file.filename.rsplit('.', 1)[0]}_nexarin.docx"
        )
        
    except Exception as e:
        print(f"Conversion error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        # We don't delete output_docx here because FileResponse needs to read it.
        # But we can delete the input_pdf.
        if os.path.exists(input_pdf):
            try:
                os.remove(input_pdf)
            except:
                pass

@app.post("/convert/img-to-pdf")
async def convert_img_to_pdf(file: UploadFile = File(...)):
    valid_extensions = ('.jpg', '.jpeg', '.png')
    if not file.filename.lower().endswith(valid_extensions):
        raise HTTPException(status_code=400, detail="File must be a JPG or PNG image")
    
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    input_img = os.path.join(temp_dir, f"input_{unique_id}_{file.filename}")
    output_pdf = os.path.join(temp_dir, f"output_{unique_id}.pdf")
    
    try:
        # Save uploaded image to temp file
        with open(input_img, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        # Open image using Pillow
        image = Image.open(input_img)
        
        # Convert image to RGB if it has alpha channel (PNG) or is in another mode
        if image.mode in ("RGBA", "P", "LA"):
            image = image.convert("RGB")
        elif image.mode != "RGB":
            image = image.convert("RGB")
            
        # Save as PDF
        image.save(output_pdf, "PDF", resolution=100.0)
        
        if not os.path.exists(output_pdf):
            raise Exception("PDF file was not generated")
            
        return FileResponse(
            output_pdf, 
            media_type="application/pdf",
            filename=f"{file.filename.rsplit('.', 1)[0]}_nexarin.pdf"
        )
        
    except Exception as e:
        print(f"Image conversion error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_img):
            try:
                os.remove(input_img)
            except:
                pass

@app.post("/convert/compress-pdf")
async def compress_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File must be a PDF")
    
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    input_pdf = os.path.join(temp_dir, f"input_comp_{unique_id}.pdf")
    output_pdf = os.path.join(temp_dir, f"output_comp_{unique_id}.pdf")
    
    try:
        # Save uploaded PDF to temp file
        with open(input_pdf, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        # Open PDF using PyMuPDF and save with compression options
        doc = fitz.open(input_pdf)
        doc.save(output_pdf, garbage=4, deflate=True, clean=True)
        doc.close()
        
        if not os.path.exists(output_pdf):
            raise Exception("Compressed PDF file was not generated")
            
        return FileResponse(
            output_pdf, 
            media_type="application/pdf",
            filename=f"{file.filename.rsplit('.', 1)[0]}_compressed.pdf"
        )
        
    except Exception as e:
        print(f"PDF compression error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_pdf):
            try:
                os.remove(input_pdf)
            except:
                pass

@app.post("/convert/pdf-to-excel")
async def convert_pdf_to_excel(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File must be a PDF")
    
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    input_pdf = os.path.join(temp_dir, f"input_excel_{unique_id}.pdf")
    output_excel = os.path.join(temp_dir, f"output_excel_{unique_id}.xlsx")
    
    try:
        with open(input_pdf, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        with pdfplumber.open(input_pdf) as pdf:
            all_tables = []
            for page in pdf.pages:
                table = page.extract_table()
                if table:
                    all_tables.extend(table)
            
            if not all_tables:
                raise Exception("Tidak ada tabel yang terdeteksi di dalam file PDF ini.")
                
            # Filter empty rows and handle columns
            cleaned_tables = [[str(cell) if cell is not None else "" for cell in row] for row in all_tables]
            
            # Use first row as header if it has data, else generate default headers
            if len(cleaned_tables) > 1:
                df = pd.DataFrame(cleaned_tables[1:], columns=cleaned_tables[0])
            else:
                df = pd.DataFrame(cleaned_tables)
                
            df.to_excel(output_excel, index=False)
            
        return FileResponse(
            output_excel, 
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=f"{file.filename.rsplit('.', 1)[0]}_nexarin.xlsx"
        )
    except Exception as e:
        print(f"PDF to Excel error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_pdf):
            try: os.remove(input_pdf)
            except: pass

@app.post("/convert/word-to-pdf")
async def convert_word_to_pdf(file: UploadFile = File(...)):
    valid_extensions = ('.doc', '.docx')
    if not file.filename.lower().endswith(valid_extensions):
        raise HTTPException(status_code=400, detail="File must be a Word document")
    
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    # Secure filename
    safe_name = f"input_{unique_id}.docx"
    input_docx = os.path.join(temp_dir, safe_name)
    expected_output_pdf = os.path.join(temp_dir, f"input_{unique_id}.pdf")
    
    try:
        with open(input_docx, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        # Using LibreOffice in headless mode
        subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, input_docx
        ], check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        if not os.path.exists(expected_output_pdf):
            raise Exception("PDF file was not generated by LibreOffice")
            
        return FileResponse(
            expected_output_pdf, 
            media_type="application/pdf",
            filename=f"{file.filename.rsplit('.', 1)[0]}_nexarin.pdf"
        )
    except Exception as e:
        print(f"Word to PDF error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_docx):
            try: os.remove(input_docx)
            except: pass

@app.post("/convert/pdf-to-png")
async def convert_pdf_to_png(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File must be a PDF")
    
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    input_pdf = os.path.join(temp_dir, f"input_png_{unique_id}.pdf")
    output_zip = os.path.join(temp_dir, f"output_pngs_{unique_id}.zip")
    
    try:
        with open(input_pdf, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        doc = fitz.open(input_pdf)
        png_files = []
        
        for i, page in enumerate(doc):
            pix = page.get_pixmap(dpi=150)
            img_path = os.path.join(temp_dir, f"page_{i+1}_{unique_id}.png")
            pix.save(img_path)
            png_files.append(img_path)
        doc.close()
        
        with zipfile.ZipFile(output_zip, 'w') as zf:
            for i, f in enumerate(png_files):
                zf.write(f, f"page_{i+1}.png")
                
        # Clean up the individual PNGs right away
        for f in png_files:
            try: os.remove(f)
            except: pass
            
        return FileResponse(
            output_zip, 
            media_type="application/zip",
            filename=f"{file.filename.rsplit('.', 1)[0]}_images.zip"
        )
    except Exception as e:
        print(f"PDF to PNG error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_pdf):
            try: os.remove(input_pdf)
            except: pass

@app.post("/convert/edit-pdf")
async def edit_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="File must be a PDF")
    
    unique_id = f"{int(time.time())}_{random.randint(1000, 9999)}"
    temp_dir = tempfile.gettempdir()
    
    input_pdf = os.path.join(temp_dir, f"input_edit_{unique_id}.pdf")
    output_pdf = os.path.join(temp_dir, f"output_edit_{unique_id}.pdf")
    
    try:
        with open(input_pdf, "wb") as buffer:
            content = await file.read()
            buffer.write(content)
            
        # Basic edit logic: Adding watermark
        # In a real app, you would accept form fields with edit commands
        doc = fitz.open(input_pdf)
        text = "EDITED WITH NEXARIN"
        
        for page in doc:
            rect = page.rect
            page.insert_text(fitz.Point(50, 50), text, fontsize=20, color=(1, 0, 0))
            
        doc.save(output_pdf)
        doc.close()
        
        return FileResponse(
            output_pdf, 
            media_type="application/pdf",
            filename=f"{file.filename.rsplit('.', 1)[0]}_edited.pdf"
        )
    except Exception as e:
        print(f"Edit PDF error: {e}")
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        if os.path.exists(input_pdf):
            try: os.remove(input_pdf)
            except: pass

# ---------------------------------------------------------
# NEW FEATURES: Merge, Split, PDF to PPTX, PPTX to PDF, Excel to PDF
# ---------------------------------------------------------

@app.post("/convert/merge-pdf")
async def merge_pdf(files: List[UploadFile] = File(...)):
    if len(files) < 2:
        raise HTTPException(status_code=400, detail="Pilih minimal 2 file PDF untuk digabungkan.")
        
    try:
        merged_doc = fitz.open()
        
        for file in files:
            if not file.filename.lower().endswith('.pdf'):
                raise HTTPException(status_code=400, detail="Semua file harus berformat PDF.")
            doc_bytes = await file.read()
            doc = fitz.open(stream=doc_bytes, filetype="pdf")
            merged_doc.insert_pdf(doc)
            doc.close()
            
        output_bytes = merged_doc.write()
        merged_doc.close()
        
        return StreamingResponse(
            io.BytesIO(output_bytes), 
            media_type="application/pdf", 
            headers={"Content-Disposition": 'attachment; filename="merged_document.pdf"'}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Terjadi kesalahan saat menggabungkan PDF: {str(e)}")

@app.post("/convert/split-pdf")
async def split_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Hanya menerima file PDF.")
        
    try:
        doc_bytes = await file.read()
        doc = fitz.open(stream=doc_bytes, filetype="pdf")
        
        zip_buffer = io.BytesIO()
        with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
            for i in range(len(doc)):
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=i, to_page=i)
                page_bytes = new_doc.write()
                new_doc.close()
                
                zip_file.writestr(f"page_{i+1}.pdf", page_bytes)
                
        doc.close()
        zip_buffer.seek(0)
        
        base_name = os.path.splitext(file.filename)[0]
        return StreamingResponse(
            zip_buffer, 
            media_type="application/zip", 
            headers={"Content-Disposition": f'attachment; filename="{base_name}_split.zip"'}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Terjadi kesalahan saat memecah PDF: {str(e)}")

@app.post("/convert/pdf-to-pptx")
async def pdf_to_pptx(file: UploadFile = File(...)):
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Hanya menerima file PDF.")
        
    try:
        doc_bytes = await file.read()
        doc = fitz.open(stream=doc_bytes, filetype="pdf")
        
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6]
        
        if len(doc) > 0:
            first_page = doc[0]
            width = first_page.rect.width
            height = first_page.rect.height
            prs.slide_width = Inches(width / 72)
            prs.slide_height = Inches(height / 72)
        
        with tempfile.TemporaryDirectory() as temp_dir:
            for i in range(len(doc)):
                page = doc[i]
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_path = os.path.join(temp_dir, f"page_{i}.png")
                pix.save(img_path)
                
                slide = prs.slides.add_slide(blank_slide_layout)
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            
        doc.close()
        pptx_buffer.seek(0)
        
        base_name = os.path.splitext(file.filename)[0]
        return StreamingResponse(
            pptx_buffer, 
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", 
            headers={"Content-Disposition": f'attachment; filename="{base_name}.pptx"'}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Terjadi kesalahan saat mengonversi PDF ke PPTX: {str(e)}")

@app.post("/convert/pptx-to-pdf")
async def pptx_to_pdf(file: UploadFile = File(...)):
    if not (file.filename.lower().endswith('.pptx') or file.filename.lower().endswith('.ppt')):
        raise HTTPException(status_code=400, detail="Hanya menerima file PPT atau PPTX.")
        
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file.filename)
            with open(input_path, "wb") as f:
                content = await file.read()
                f.write(content)
                
            subprocess.run([
                "libreoffice", "--headless", "--invisible",
                "--convert-to", "pdf",
                "--outdir", temp_dir,
                input_path
            ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            
            base_name = os.path.splitext(file.filename)[0]
            pdf_filename = f"{base_name}.pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)
            
            if not os.path.exists(pdf_path):
                raise Exception("File PDF gagal dibuat oleh LibreOffice.")
                
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
                
            return StreamingResponse(
                io.BytesIO(pdf_bytes),
                media_type="application/pdf",
                headers={"Content-Disposition": f'attachment; filename="{pdf_filename}"'}
            )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Terjadi kesalahan saat mengonversi PPTX ke PDF: {str(e)}")

@app.post("/convert/excel-to-pdf")
async def excel_to_pdf(file: UploadFile = File(...)):
    if not (file.filename.lower().endswith('.xlsx') or file.filename.lower().endswith('.xls') or file.filename.lower().endswith('.csv')):
        raise HTTPException(status_code=400, detail="Hanya menerima file Excel (XLS/XLSX/CSV).")
        
    try:
        with tempfile.TemporaryDirectory() as temp_dir:
            input_path = os.path.join(temp_dir, file.filename)
            with open(input_path, "wb") as f:
                content = await file.read()
                f.write(content)
                
            subprocess.run([
                "libreoffice", "--headless", "--invisible",
                "--convert-to", "pdf",
                "--outdir", temp_dir,
                input_path
            ], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            
            base_name = os.path.splitext(file.filename)[0]
            pdf_filename = f"{base_name}.pdf"
            pdf_path = os.path.join(temp_dir, pdf_filename)
            
            if not os.path.exists(pdf_path):
                raise Exception("File PDF gagal dibuat oleh LibreOffice.")
                
            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
                
            return StreamingResponse(
                io.BytesIO(pdf_bytes),
                media_type="application/pdf",
                headers={"Content-Disposition": f'attachment; filename="{pdf_filename}"'}
            )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Terjadi kesalahan saat mengonversi Excel ke PDF: {str(e)}")
